"""
문서 → 청크 → ChromaDB 적재 (배치 처리)

BATCH_SIZE 개 파일씩 추출 → 청킹 → 임베딩 → ChromaDB 기록
체크포인트: data/chunks/checkpoint.json (--resume 로 이어서 실행)

실행:  python -m src.retrievers.ingest
재시작: python -m src.retrievers.ingest --resume
초기화: python -m src.retrievers.ingest --reset
"""
import gc
import json
import re
import sys
import time
import zipfile
import shutil
from pathlib import Path

import chromadb
from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_ollama import OllamaEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter

from config.settings import settings

BATCH_SIZE = 50
CHROMA_MAX = 5000  # chromadb 단일 upsert 한도(5461)보다 여유 있게

AUTO_DELETE_EXTS = {
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif", ".svg", ".ico",
    ".mp4", ".avi", ".mov", ".mkv", ".wmv", ".flv", ".webm", ".m4v",
    ".mp3", ".wav", ".aac", ".flac", ".m4a", ".ogg",
    ".ttf", ".otf", ".woff", ".woff2", ".eot",
}

CHECKPOINT_FILE = Path(settings.chunks_dir) / "checkpoint.json"

SUMMARY_SPLITTER = RecursiveCharacterTextSplitter(
    chunk_size=2000, chunk_overlap=400,
    separators=["\n\n", "\n", "。", ". ", " "],
)
DETAIL_SPLITTER = RecursiveCharacterTextSplitter(
    chunk_size=500, chunk_overlap=100,
    separators=["\n\n", "\n", "。", ". ", " "],
)


# ── 텍스트 추출 ─────────────────────────────────────────────────────────────

def _extract_pdf(path: Path) -> str:
    import pdfplumber
    texts = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                texts.append(t)
    return "\n\n".join(texts)


def _extract_ppt_legacy(path: Path) -> str:
    """구형 .ppt OLE 바이너리에서 텍스트 레코드 직접 추출.

    PPT 레코드는 중첩 컨테이너 구조(recVer==0xF)이므로
    컨테이너는 헤더만 소비하고 내부를 계속 순회한다.
    """
    import olefile, struct
    texts = []
    try:
        if not olefile.isOleFile(str(path)):
            return ""
        with olefile.OleFileIO(str(path)) as ole:
            if not ole.exists("PowerPoint Document"):
                return ""
            data = ole.openstream("PowerPoint Document").read()

        i = 0
        while i + 8 <= len(data):
            word0    = struct.unpack_from("<H", data, i)[0]
            rec_ver  = word0 & 0x000F          # 하위 4비트 = 버전
            rec_type = struct.unpack_from("<H", data, i + 2)[0]
            rec_len  = struct.unpack_from("<I", data, i + 4)[0]
            i += 8  # 헤더 소비

            if rec_ver == 0x0F:
                # 컨테이너 레코드 — 내부를 바로 순회
                continue

            if i + rec_len > len(data):
                break

            chunk = data[i: i + rec_len]

            if rec_type == 0x0FA0 and rec_len > 0:   # TextCharsAtom (UTF-16LE)
                txt = chunk.decode("utf-16-le", errors="ignore").strip()
                if txt:
                    texts.append(txt)
            elif rec_type == 0x0FA8 and rec_len > 0:  # TextBytesAtom (Latin-1)
                txt = chunk.decode("latin-1", errors="ignore").strip()
                if txt:
                    texts.append(txt)

            i += rec_len
    except Exception:
        pass
    return "\n".join(texts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    try:
        prs = Presentation(str(path))
    except Exception:
        # 구형 .ppt OLE 바이너리 폴백
        return _extract_ppt_legacy(path)
    texts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
        if lines:
            texts.append(f"[슬라이드 {slide_num}]\n" + "\n".join(lines))
    return "\n\n".join(texts)


def _extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    texts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            texts.append(f"[시트: {sheet.title}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(texts)


def _extract_xls(path: Path) -> str:
    import xlrd
    wb = xlrd.open_workbook(str(path))
    texts = []
    for sheet in wb.sheets():
        rows = []
        for r in range(sheet.nrows):
            cells = [str(sheet.cell_value(r, c)).strip()
                     for c in range(sheet.ncols)
                     if str(sheet.cell_value(r, c)).strip()]
            if cells:
                rows.append("\t".join(cells))
        if rows:
            texts.append(f"[시트: {sheet.name}]\n" + "\n".join(rows))
    return "\n\n".join(texts)


def _extract_docx(path: Path) -> str:
    from docx import Document as DocxDoc
    doc = DocxDoc(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                paras.append("\t".join(cells))
    return "\n\n".join(paras)


def _extract_hwpx(path: Path) -> str:
    texts = []
    with zipfile.ZipFile(str(path), "r") as z:
        for name in z.namelist():
            if name.startswith("Contents/") and name.endswith(".xml"):
                xml = z.read(name).decode("utf-8", errors="ignore")
                texts.extend(re.findall(r"<t[^>]*>([^<]+)</t>", xml))
    return "\n".join(t.strip() for t in texts if t.strip())


def _extract_hwp(path: Path) -> str:
    import olefile, struct
    texts = []
    try:
        with olefile.OleFileIO(str(path)) as ole:
            for entry in ole.listdir():
                if entry[0] == "BodyText":
                    data = ole.openstream(entry).read()
                    try:
                        import zlib
                        data = zlib.decompress(data, -15)
                    except Exception:
                        pass
                    i = 0
                    while i + 4 <= len(data):
                        tag_id = struct.unpack_from("<H", data, i)[0] & 0x3FF
                        size   = struct.unpack_from("<H", data, i + 2)[0]
                        i += 4
                        if i + size > len(data):
                            break
                        if tag_id == 67:
                            txt = data[i:i+size].decode("utf-16-le", errors="ignore").strip()
                            if txt:
                                texts.append(txt)
                        i += size
    except Exception:
        pass
    if texts:
        return "\n".join(texts)
    raw = path.read_bytes()
    decoded = raw.decode("utf-16-le", errors="ignore")
    return "\n".join(re.findall(r"[가-힣A-Za-z0-9,. ]{5,}", decoded))


def _extract_zip(path: Path) -> str:
    extract_dir = path.parent / f"__unzip_{path.stem}__"
    texts = []
    try:
        with zipfile.ZipFile(str(path), "r") as z:
            z.extractall(str(extract_dir))
        for p in extract_dir.rglob("*"):
            if p.is_file() and p.suffix.lower() in EXTRACTORS:
                try:
                    t = EXTRACTORS[p.suffix.lower()](p)
                    if t.strip():
                        texts.append(t)
                except Exception:
                    pass
    except Exception:
        pass
    finally:
        if extract_dir.exists():
            shutil.rmtree(extract_dir)
    return "\n\n".join(texts)


def _extract_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


EXTRACTORS = {
    ".pdf": _extract_pdf, ".pptx": _extract_pptx, ".ppt": _extract_pptx,
    ".xlsx": _extract_xlsx, ".xls": _extract_xls, ".docx": _extract_docx,
    ".hwpx": _extract_hwpx, ".hwp": _extract_hwp, ".zip": _extract_zip,
    ".txt": _extract_text, ".md": _extract_text,
}


# ── 메타데이터 ───────────────────────────────────────────────────────────────

def _meta(text: str, source: str, source_type: str = "document") -> dict:
    return {
        "source": source,
        "source_type": source_type,
        "has_si":         bool(re.search(r"시스템.{0,10}구축|SI|정보화", text)),
        "has_consulting": bool(re.search(r"컨설팅|용역|자문", text)),
        "has_security":   bool(re.search(r"보안|정보보호|취약점", text)),
        "has_rfp":        bool(re.search(r"제안요청서|RFP|과업지시서", text)),
        "has_data":       bool(re.search(r"데이터|빅데이터|AI|인공지능", text)),
    }


# ── 체크포인트 ───────────────────────────────────────────────────────────────

def load_checkpoint() -> set:
    if CHECKPOINT_FILE.exists():
        return set(json.loads(CHECKPOINT_FILE.read_text(encoding="utf-8")))
    return set()


def save_checkpoint(done: set):
    CHECKPOINT_FILE.parent.mkdir(parents=True, exist_ok=True)
    CHECKPOINT_FILE.write_text(
        json.dumps(sorted(done), ensure_ascii=False), encoding="utf-8"
    )


# ── 배치 처리 ────────────────────────────────────────────────────────────────

def process_batch(
    paths: list[Path],
    summary_store: Chroma,
    detail_store: Chroma,
) -> tuple[list[str], list[str]]:
    """파일 목록을 추출 → 청킹 → ChromaDB 적재. (성공명, 실패명) 반환."""
    summary_chunks, detail_chunks = [], []
    succeeded, failed = [], []

    for path in paths:
        ext = path.suffix.lower()
        if ext not in EXTRACTORS:
            continue
        try:
            text = EXTRACTORS[ext](path)
            if not text.strip():
                print(f"  ⚠  {path.name} (내용 없음, 스킵)", flush=True)
                succeeded.append(path.name)
                continue

            doc = Document(page_content=text, metadata=_meta(text, path.name))

            for chunk in SUMMARY_SPLITTER.split_documents([doc]):
                chunk.metadata["chunk_type"] = "summary"
                summary_chunks.append(chunk)
            for chunk in DETAIL_SPLITTER.split_documents([doc]):
                chunk.metadata["chunk_type"] = "detail"
                detail_chunks.append(chunk)

            print(f"  ✓ {path.name} ({len(text):,}자)", flush=True)
            succeeded.append(path.name)
        except Exception as e:
            print(f"  ✗ {path.name} ({type(e).__name__}: {e})", flush=True)
            failed.append(path.name)

    for label, store, chunks in [
        ("요약", summary_store, summary_chunks),
        ("세부", detail_store,  detail_chunks),
    ]:
        if not chunks:
            continue
        for start in range(0, len(chunks), CHROMA_MAX):
            sub = chunks[start: start + CHROMA_MAX]
            print(
                f"  → {label} {start+1}~{start+len(sub)}/{len(chunks)}청크 적재 중...",
                flush=True,
            )
            store.add_documents(sub)

    del summary_chunks, detail_chunks
    gc.collect()

    return succeeded, failed


# ── 메인 ─────────────────────────────────────────────────────────────────────

def main():
    resume = "--resume" in sys.argv
    reset  = "--reset"  in sys.argv

    raw_dir = Path(settings.raw_data_dir)

    # 지원 파일 목록 (자동 삭제 대상 제외)
    all_paths = sorted([
        p for p in raw_dir.rglob("*")
        if p.is_file()
        and p.suffix.lower() in EXTRACTORS
        and p.suffix.lower() not in AUTO_DELETE_EXTS
    ])
    total = len(all_paths)
    print(f"=== ChromaDB 배치 적재 시작 ===", flush=True)
    print(f"대상 파일: {total}개 | 배치 크기: {BATCH_SIZE}", flush=True)

    # 체크포인트 처리
    done_set: set[str] = set()
    if reset:
        if CHECKPOINT_FILE.exists():
            CHECKPOINT_FILE.unlink()
        chroma_client = chromadb.PersistentClient(path=settings.chroma_persist_dir)
        for col in [settings.summary_collection, settings.detail_collection]:
            try:
                chroma_client.delete_collection(col)
                print(f"  기존 컬렉션 삭제: {col}", flush=True)
            except Exception:
                pass
    elif resume:
        done_set = load_checkpoint()
        print(f"이어서 실행: 이미 완료 {len(done_set)}개", flush=True)

    # 남은 파일
    pending = [p for p in all_paths if p.name not in done_set]
    print(f"처리 예정: {len(pending)}개\n", flush=True)

    if not pending:
        print("모든 파일이 이미 처리되었습니다.", flush=True)
        return

    # ChromaDB 스토어 (get_or_create)
    chroma_client = chromadb.PersistentClient(path=settings.chroma_persist_dir)
    embedding = OllamaEmbeddings(model=settings.exaone_model)

    summary_store = Chroma(
        client=chroma_client,
        collection_name=settings.summary_collection,
        embedding_function=embedding,
    )
    detail_store = Chroma(
        client=chroma_client,
        collection_name=settings.detail_collection,
        embedding_function=embedding,
    )

    # 배치 루프
    all_failed = []
    t0 = time.time()

    for batch_idx in range(0, len(pending), BATCH_SIZE):
        batch = pending[batch_idx: batch_idx + BATCH_SIZE]
        batch_num = batch_idx // BATCH_SIZE + 1
        total_batches = (len(pending) + BATCH_SIZE - 1) // BATCH_SIZE

        elapsed = time.time() - t0
        done_so_far = len(done_set)
        print(
            f"\n[배치 {batch_num}/{total_batches}] "
            f"파일 {batch_idx+1}~{min(batch_idx+BATCH_SIZE, len(pending))} / {len(pending)} "
            f"| 경과 {elapsed/60:.1f}분",
            flush=True,
        )

        succeeded, failed = process_batch(batch, summary_store, detail_store)
        done_set.update(succeeded)
        all_failed.extend(failed)

        save_checkpoint(done_set)
        print(
            f"  배치 완료 ✓ | 누적 {len(done_set)}/{total} | 실패 {len(all_failed)}",
            flush=True,
        )

    total_elapsed = (time.time() - t0) / 60
    print(f"\n=== 적재 완료 ===", flush=True)
    print(f"성공: {len(done_set)}개 | 실패: {len(all_failed)}개 | 소요: {total_elapsed:.1f}분", flush=True)

    if all_failed:
        print("\n실패 목록:", flush=True)
        for f in all_failed:
            print(f"  - {f}", flush=True)


def ingest_single_file(path: Path) -> dict:
    """단일 파일을 추출 → 청킹 → ChromaDB 적재. 결과 dict 반환."""
    ext = path.suffix.lower()
    if ext not in EXTRACTORS:
        return {"success": False, "error": f"지원하지 않는 파일 형식: {ext}"}

    try:
        text = EXTRACTORS[ext](path)
    except Exception as e:
        return {"success": False, "error": f"텍스트 추출 실패: {e}"}

    if not text.strip():
        return {"success": False, "error": "텍스트 내용 없음 (이미지 전용 문서일 수 있음)"}

    chroma_client = chromadb.PersistentClient(path=settings.chroma_persist_dir)
    embedding = OllamaEmbeddings(model=settings.exaone_model)
    summary_store = Chroma(
        client=chroma_client,
        collection_name=settings.summary_collection,
        embedding_function=embedding,
    )
    detail_store = Chroma(
        client=chroma_client,
        collection_name=settings.detail_collection,
        embedding_function=embedding,
    )

    doc = Document(page_content=text, metadata=_meta(text, path.name))

    summary_chunks = []
    for chunk in SUMMARY_SPLITTER.split_documents([doc]):
        chunk.metadata["chunk_type"] = "summary"
        summary_chunks.append(chunk)

    detail_chunks = []
    for chunk in DETAIL_SPLITTER.split_documents([doc]):
        chunk.metadata["chunk_type"] = "detail"
        detail_chunks.append(chunk)

    for store, chunks in [(summary_store, summary_chunks), (detail_store, detail_chunks)]:
        for start in range(0, len(chunks), CHROMA_MAX):
            store.add_documents(chunks[start: start + CHROMA_MAX])

    done = load_checkpoint()
    done.add(path.name)
    save_checkpoint(done)

    del summary_chunks, detail_chunks
    gc.collect()

    return {
        "success": True,
        "filename": path.name,
        "chars": len(text),
    }


def get_file_chunks(filename: str) -> dict:
    """ChromaDB에서 특정 파일의 청크 수와 첫 번째 요약 청크 내용을 반환."""
    try:
        client = chromadb.PersistentClient(path=settings.chroma_persist_dir)
        summary_col = client.get_collection(settings.summary_collection)
        detail_col  = client.get_collection(settings.detail_collection)

        s_result = summary_col.get(where={"source": filename}, include=["documents", "metadatas"])
        d_result = detail_col.get(where={"source": filename}, include=["documents"])

        preview = ""
        if s_result["documents"]:
            preview = s_result["documents"][0][:1200]

        return {
            "summary_chunks": len(s_result["documents"]),
            "detail_chunks": len(d_result["documents"]),
            "preview": preview,
        }
    except Exception as e:
        return {"summary_chunks": 0, "detail_chunks": 0, "preview": "", "error": str(e)}


def ingest_note(filename: str, note_text: str) -> dict:
    """파일에 대한 보충 노트를 ChromaDB에 적재. 기존 노트 청크는 삭제 후 재적재."""
    if not note_text.strip():
        return {"success": True, "message": "노트 내용 없음"}

    note_source = f"__note__{filename}"
    client = chromadb.PersistentClient(path=settings.chroma_persist_dir)
    embedding = OllamaEmbeddings(model=settings.exaone_model)

    summary_store = Chroma(
        client=client, collection_name=settings.summary_collection,
        embedding_function=embedding,
    )
    detail_store = Chroma(
        client=client, collection_name=settings.detail_collection,
        embedding_function=embedding,
    )

    # 기존 노트 청크 삭제
    for store in (summary_store, detail_store):
        try:
            existing = store.get(where={"source": note_source})
            if existing["ids"]:
                store.delete(ids=existing["ids"])
        except Exception:
            pass

    doc = Document(
        page_content=note_text,
        metadata=_meta(note_text, note_source, source_type="note"),
    )

    summary_chunks = []
    for chunk in SUMMARY_SPLITTER.split_documents([doc]):
        chunk.metadata["chunk_type"] = "summary"
        summary_chunks.append(chunk)
    detail_chunks = []
    for chunk in DETAIL_SPLITTER.split_documents([doc]):
        chunk.metadata["chunk_type"] = "detail"
        detail_chunks.append(chunk)

    for store, chunks in [(summary_store, summary_chunks), (detail_store, detail_chunks)]:
        if chunks:
            store.add_documents(chunks)

    return {"success": True, "chunks": len(summary_chunks)}


def get_similar_docs(query: str, k: int = 5) -> list[dict]:
    """진단용: 쿼리에 대한 상위 문서와 L2 거리를 반환."""
    try:
        client = chromadb.PersistentClient(path=settings.chroma_persist_dir)
        embedding = OllamaEmbeddings(model=settings.exaone_model)
        store = Chroma(
            client=client, collection_name=settings.summary_collection,
            embedding_function=embedding,
        )
        results = store.similarity_search_with_score(query, k=k)
        return [
            {
                "source": r.metadata.get("source", ""),
                "distance": round(score, 4),
                "preview": r.page_content[:200],
            }
            for r, score in results
        ]
    except Exception as e:
        return [{"error": str(e)}]


if __name__ == "__main__":
    main()
